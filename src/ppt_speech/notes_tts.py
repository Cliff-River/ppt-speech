"""PPT 备注文字转语音模块。

使用 edge-tts 将 PowerPoint 幻灯片备注转换为语音，
并嵌入到 PPT 中实现幻灯片播放时自动配音。
"""

from __future__ import annotations

import asyncio
import re
import shutil
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from edge_tts import Communicate
from edge_tts.exceptions import (
    EdgeTTSException,
    NoAudioReceived,
    UnexpectedResponse,
    UnknownResponse,
    WebSocketError,
)
from lxml import etree
from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Inches

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"

_RATE_PATTERN = re.compile(r"^[+-]\d+%$")
_VOICE_PATTERN = re.compile(r"^([a-z]{2,})-([A-Z]{2,})-(.+Neural)$")


@dataclass(slots=True)
class PTSpeechConfig:
    """PPT 配音处理配置。

    Attributes:
        input_dir: 输入目录路径，存放原始 PPT 文件。
        output_dir: 输出目录路径，存放生成的配音 PPT 文件。
        input_filename: 输入 PPT 文件名（不含目录路径）。
        output_filename: 输出 PPT 文件名（不含目录路径）。
        voice_name: 语音名称，格式为 'zh-CN-XiaoxiaoNeural'。
        speech_rate: 语速调整，格式为 '+0%' 或 '-50%'。
        temp_audio_dir: 临时音频文件存放目录。
        audio_icon_offset: 音频图标在画布上的偏移英寸数。
        audio_icon_size: 音频图标尺寸（英寸）。
    """

    input_dir: Path = field(default_factory=lambda: Path("data"))
    output_dir: Path = field(default_factory=lambda: Path("data"))
    input_filename: str = "input.pptx"
    output_filename: str = "output.pptx"
    voice_name: str = "zh-CN-XiaoxiaoNeural"
    speech_rate: str = "+0%"
    temp_audio_dir: Path = field(default_factory=lambda: Path(".tmp_audio"))
    audio_icon_offset: float = -2.0
    audio_icon_size: float = 1.0

    @property
    def input_path(self) -> Path:
        """完整输入文件路径。"""
        return self.input_dir / self.input_filename

    @property
    def output_path(self) -> Path:
        """完整输出文件路径。"""
        return self.output_dir / self.output_filename

    def validate(self) -> None:
        """验证配置参数的合法性。

        Raises:
            ValueError: 当语音名称或语速格式不正确时。
            FileNotFoundError: 当输入 PPT 文件不存在时。
        """
        if not _VOICE_PATTERN.match(self.voice_name):
            raise ValueError(
                f"语音名称格式错误: '{self.voice_name}'，"
                f"正确格式如 'zh-CN-XiaoxiaoNeural'"
            )
        if not _RATE_PATTERN.match(self.speech_rate):
            raise ValueError(
                f"语速格式错误: '{self.speech_rate}'，"
                f"正确格式如 '+0%' 或 '-50%'"
            )
        if not self.input_path.exists():
            raise FileNotFoundError(
                f"输入 PPT 文件不存在: {self.input_path}"
            )


def _normalize_voice_name(voice_name: str) -> str:
    """将简短语音名称转换为 edge-tts 完整格式。

    Args:
        voice_name: 简短语音名称，如 'zh-CN-XiaoxiaoNeural'。

    Returns:
        完整格式的语音名称，如
        'Microsoft Server Speech Text to Speech Voice (zh-CN, XiaoxiaoNeural)'。
    """
    match = _VOICE_PATTERN.match(voice_name)
    if match is None:
        return voice_name

    lang = match.group(1)
    region = match.group(2)
    name = match.group(3)
    if "-" in name:
        region = f"{region}-{name[:name.find('-')]}"
        name = name[name.find("-") + 1:]

    return (
        "Microsoft Server Speech Text to Speech Voice"
        f" ({lang}-{region}, {name})"
    )


def _read_notes_text(slide: Slide) -> str:
    """从幻灯片中提取备注文字。

    Args:
        slide: PowerPoint 幻灯片对象。

    Returns:
        备注文字内容，若无备注则返回空字符串。
    """
    if not slide.has_notes_slide:
        return ""
    notes_slide = slide.notes_slide
    if notes_slide.notes_text_frame is None:
        return ""
    return notes_slide.notes_text_frame.text.strip()


async def text_to_mp3(
    text: str,
    save_path: Path,
    voice_name: str = "zh-CN-XiaoxiaoNeural",
    speech_rate: str = "+0%",
) -> bool:
    """使用 edge-tts 将文字转换为 MP3 音频。

    Args:
        text: 要转换的文字内容。
        save_path: MP3 文件保存路径。
        voice_name: 语音名称，如 'zh-CN-XiaoxiaoNeural'。
        speech_rate: 语速调整，如 '+0%'、'-50%'。

    Returns:
        转换成功返回 True，文字为空时返回 False。

    Raises:
        ValueError: 当语速格式不正确时。
        EdgeTTSException: 当 TTS 服务请求失败时。
        OSError: 当文件保存失败时。
    """
    if not text.strip():
        return False

    if not _RATE_PATTERN.match(speech_rate):
        raise ValueError(f"语速格式错误: '{speech_rate}'，应为 '+0%' 格式")

    full_voice = _normalize_voice_name(voice_name)
    communicate = Communicate(text, full_voice, rate=speech_rate)

    try:
        save_path.parent.mkdir(parents=True, exist_ok=True)
        await communicate.save(str(save_path))
    except (NoAudioReceived, UnexpectedResponse, UnknownResponse, WebSocketError) as exc:
        raise EdgeTTSException(f"TTS 合成失败: {exc}") from exc
    except OSError as exc:
        raise OSError(f"音频文件保存失败 '{save_path}': {exc}") from exc

    return True


def _apply_autoplay_timing(slide_element) -> None:
    """修改幻灯片 XML 实现音频自动播放。

    通过设置 playOnEntry 属性和时序配置，
    使幻灯片进入时音频自动播放。

    Args:
        slide_element: 幻灯片的 lxml 元素对象。
    """
    for media_el in slide_element.iter(f"{{{P14_NS}}}media"):
        media_el.set(f"{{{P_NS}}}playOnEntry", "1")

    timing = slide_element.find(f"{{{P_NS}}}timing")
    if timing is None:
        return

    for cond in timing.findall(f".//{{{P_NS}}}cond"):
        cond.set("delay", "0")
        if cond.find(f"{{{P_NS}}}withPrev") is None:
            etree.SubElement(cond, f"{{{P_NS}}}withPrev")

    for media_node in timing.findall(f".//{{{P_NS}}}cMediaNode"):
        media_node.set(f"{{{P_NS}}}playOnEntry", "1")


def embed_audio_autoplay(
    slide: Slide,
    audio_path: Path,
    icon_offset: float = -2.0,
    icon_size: float = 1.0,
) -> None:
    """将音频嵌入幻灯片并设置为自动播放。

    音频图标放置在画布外以实现视觉隐藏，
    通过修改 XML 实现幻灯片进入时自动播放。

    Args:
        slide: PowerPoint 幻灯片对象。
        audio_path: 音频文件路径。
        icon_offset: 图标在画布上的偏移英寸数（负值表示画布外）。
        icon_size: 图标尺寸（英寸）。

    Raises:
        FileNotFoundError: 当音频文件不存在时。
        ValueError: 当图标尺寸或偏移为负数时。
    """
    if not audio_path.exists():
        raise FileNotFoundError(f"音频文件不存在: {audio_path}")

    if icon_size <= 0:
        raise ValueError(f"图标尺寸必须为正数，当前值: {icon_size}")

    slide.shapes.add_movie(
        str(audio_path),
        left=Inches(icon_offset),
        top=Inches(icon_offset),
        width=Inches(icon_size),
        height=Inches(icon_size),
        poster_frame_image=None,
        mime_type="audio/mpeg",
    )

    _apply_autoplay_timing(slide._element)


async def process_slides(
    prs: Presentation,
    config: PTSpeechConfig,
) -> None:
    """处理演示文稿中的所有幻灯片，生成并嵌入配音。

    Args:
        prs: PowerPoint 演示文稿对象。
        config: 配音处理配置。

    Raises:
        EdgeTTSException: 当 TTS 合成失败时。
        OSError: 当输出文件保存失败时。
    """
    config.temp_audio_dir.mkdir(parents=True, exist_ok=True)

    try:
        for idx, slide in enumerate(prs.slides, start=1):
            note_text = _read_notes_text(slide)

            if not note_text:
                print(f"【第{idx}页】无备注，跳过配音")
                continue

            preview = note_text[:30] + "..." if len(note_text) > 30 else note_text
            print(f"【第{idx}页】生成语音：{preview}")

            audio_file = config.temp_audio_dir / f"slide_{idx}.mp3"
            success = await text_to_mp3(
                note_text,
                audio_file,
                voice_name=config.voice_name,
                speech_rate=config.speech_rate,
            )

            if success:
                embed_audio_autoplay(
                    slide,
                    audio_file,
                    icon_offset=config.audio_icon_offset,
                    icon_size=config.audio_icon_size,
                )

        config.output_dir.mkdir(parents=True, exist_ok=True)
        prs.save(str(config.output_path))
        print(f"\n✅ 处理完成！输出文件：{config.output_path}")

    finally:
        if config.temp_audio_dir.exists():
            shutil.rmtree(config.temp_audio_dir, ignore_errors=True)


async def main(config: Optional[PTSpeechConfig] = None) -> None:
    """PPT 配音处理主入口。

    Args:
        config: 配音处理配置，为 None 时使用默认配置。

    Raises:
        ValueError: 当配置参数不合法时。
        FileNotFoundError: 当输入文件不存在时。
        EdgeTTSException: 当 TTS 合成失败时。
        OSError: 当文件操作失败时。
    """
    if config is None:
        config = PTSpeechConfig()

    config.validate()

    prs = Presentation(str(config.input_path))
    await process_slides(prs, config)


if __name__ == "__main__":
    asyncio.run(main())