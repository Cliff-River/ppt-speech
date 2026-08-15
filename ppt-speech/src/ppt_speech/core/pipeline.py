"""PPT 配音流程编排模块（core 子包）。

本模块是 PPT 语音合成功能的顶层入口，负责按顺序编排以下步骤：

1. 验证配置（使用 :mod:`ppt_speech.core.config`）
2. 打开 PPT 演示文稿
3. 逐页提取备注（使用 :mod:`ppt_speech.core.notes_reader`）
4. 调用 TTS 服务生成音频（使用 :mod:`ppt_speech.core.tts_client`）
5. 将音频嵌入幻灯片并配置自动播放（使用 :mod:`ppt_speech.core.audio`）
6. 按需按「音频时长 + n 秒」设置自动翻页（使用 :mod:`ppt_speech.core.slide_transition`）
7. 保存输出 PPT 并清理临时音频文件

中间音频文件统一存放在临时目录中：默认通过 :func:`tempfile.TemporaryDirectory`
在系统临时目录下创建，避免在当前工作目录留下硬编码路径；处理结束（无论成功
与否）后由上下文管理器自动清理。

进度回调
========
``process_slides`` / ``speak_ppt_notes`` 支持可选的 ``on_progress`` 回调，
用于将处理进度以结构化事件形式上报（阶段、页码、百分比、预计剩余时间）。
此机制使同一套编排逻辑既能服务 CLI（回调为 None 时走原 ``print`` 输出），
又能服务服务端（回调将事件写入 Redis 并经 SSE 推送）。本模块**不引入任何
外部依赖**，回调协议仅依赖标准库类型。

本模块仅承担「编排」职责，对外暴露 :func:`speak_ppt_notes` 与
:func:`process_slides` 两个入口；其余子模块的公共符号统一由 :mod:`ppt_speech.core`
子包重新导出。
"""

from __future__ import annotations

import shutil
import tempfile
import time
from contextlib import contextmanager
from pathlib import Path
from typing import Callable, Optional, TypedDict

from pptx import Presentation
from pptx.slide import Slide

from ppt_speech.core.audio import embed_audio_autoplay, get_audio_duration
from ppt_speech.core.config import PTSpeechConfig
from ppt_speech.core.notes_reader import read_notes_text
from ppt_speech.core.slide_transition import set_advance_after_time
from ppt_speech.core.tts_client import text_to_mp3

__all__ = ["speak_ppt_notes", "process_slides"]


# ---------------------------------------------------------------------------
# 进度回调协议（仅依赖标准库，保持核心库零外部依赖）
# ---------------------------------------------------------------------------

# 处理阶段常量。服务端（ppt_speech.server.progress.Stage）复用同一组字符串
# 语义，并额外补充 FAILED 等服务端专属状态。
STAGE_VALIDATING = "VALIDATING"
STAGE_READING_NOTES = "READING_NOTES"
STAGE_SYNTHESIZING = "SYNTHESIZING"
STAGE_EMBEDDING = "EMBEDDING"
STAGE_SETTING_TRANSITION = "SETTING_TRANSITION"
STAGE_SAVING = "SAVING"
STAGE_COMPLETED = "COMPLETED"


class ProcessProgressEvent(TypedDict, total=False):
    """处理进度事件结构。

    Attributes:
        stage: 当前处理阶段（见 ``STAGE_*`` 常量）。
        slide_idx: 当前幻灯片页码（1-based）；非页级阶段为 0。
        total_slides: 幻灯片总页数。
        percent: 完成百分比（0.0 ~ 100.0）。
        eta_seconds: 预计剩余秒数；早期或未知时为 None。
        message: 人类可读的进度说明。
    """

    stage: str
    slide_idx: int
    total_slides: int
    percent: float
    eta_seconds: "float | None"
    message: str


ProgressCallback = Callable[[ProcessProgressEvent], None]


# 每页内各子阶段的完成权重（页内合计 1.0），用于计算百分比。
_STAGE_WEIGHT = {
    STAGE_READING_NOTES: 0.0,
    STAGE_SYNTHESIZING: 0.5,
    STAGE_EMBEDDING: 0.75,
    STAGE_SETTING_TRANSITION: 0.9,
}


def _percent(idx: int, total: int, stage: str) -> float:
    """根据页码、总页数与阶段计算完成百分比。"""
    if stage == STAGE_COMPLETED:
        return 100.0
    if stage == STAGE_VALIDATING:
        return 0.0
    if stage == STAGE_SAVING:
        return 99.0
    weight = _STAGE_WEIGHT.get(stage, 0.0)
    if total <= 0:
        return 0.0
    return round(((idx - 1) + weight) / total * 100, 1)


def _eta_seconds(start: float, percent: float) -> "float | None":
    """根据已用时间与百分比估算剩余秒数。

    百分比过低（≤5）时返回 None 以避免早期抖动放大。
    """
    if percent <= 5:
        return None
    elapsed = time.monotonic() - start
    return round(elapsed * (100 - percent) / percent, 1)


def _emit_progress(
    on_progress: "ProgressCallback | None",
    start: float,
    stage: str,
    idx: int,
    total: int,
    message: str = "",
) -> None:
    """有回调时构造并发射一个结构化进度事件；无回调时为 no-op。"""
    if on_progress is None:
        return
    percent = _percent(idx, total, stage)
    on_progress(
        {
            "stage": stage,
            "slide_idx": idx,
            "total_slides": total,
            "percent": percent,
            "eta_seconds": _eta_seconds(start, percent),
            "message": message,
        }
    )


def _log(on_progress: "ProgressCallback | None", message: str) -> None:
    """无回调时打印消息（保持 CLI 原行为）；有回调时静默（服务端安静）。"""
    if on_progress is None:
        try:
            print(message)
        except UnicodeEncodeError:
            import sys

            encoding = getattr(sys.stdout, "encoding", None) or "utf-8"
            safe_message = message.encode(encoding, errors="replace").decode(
                encoding, errors="replace"
            )
            print(safe_message)


@contextmanager
def _temp_audio_workspace(config: PTSpeechConfig):
    """为中间音频文件创建临时工作目录，并在退出时自动清理。

    目录选取策略：

    - 若 ``config.temp_audio_dir`` 显式指定，则沿用该路径（便于调用方固定
      临时目录位置或调试），退出时由本函数负责删除。
    - 若为 ``None``（默认），则通过 :func:`tempfile.TemporaryDirectory`
      在系统临时目录中创建，避免在当前工作目录留下 ``.tmp_audio`` 等
      硬编码路径，从而提升可移植性并减少权限/依赖问题。

    无论处理成功还是抛出异常，上下文退出时都会清理临时目录，
    因此调用方无需手动删除临时文件。

    Args:
        config: 配音处理配置对象，用于决定临时目录位置。

    Yields:
        已创建、可写入的临时目录 :class:`~pathlib.Path`。

    Raises:
        OSError: 当临时目录无法创建时（例如磁盘已满或权限不足）。
    """
    temp_dir = config.temp_audio_dir

    if temp_dir is not None:
        # 调用方显式指定了临时目录：创建并在退出时负责清理。
        try:
            temp_dir.mkdir(parents=True, exist_ok=True)
        except OSError as exc:
            raise OSError(f"无法创建临时音频目录 '{temp_dir}': {exc}") from exc
        try:
            yield temp_dir
        finally:
            # 即使后续处理抛出异常，也确保临时目录被清理
            if temp_dir.exists():
                shutil.rmtree(temp_dir, ignore_errors=True)
    else:
        # 未指定临时目录：使用系统临时目录，由 TemporaryDirectory 自动清理。
        # ignore_cleanup_errors=True 保证清理过程中的非致命错误（如文件占用）
        # 不会掩盖主流程中真正重要的异常。
        try:
            workspace = tempfile.TemporaryDirectory(
                prefix="ppt_speech_", ignore_cleanup_errors=True
            )
        except OSError as exc:
            raise OSError(f"无法创建系统临时目录: {exc}") from exc
        with workspace:
            yield Path(workspace.name)


def _apply_auto_advance(
    slide: Slide,
    audio_file: Path,
    config: PTSpeechConfig,
    idx: int,
    *,
    total: int = 0,
    start: float = 0.0,
    on_progress: "ProgressCallback | None" = None,
) -> None:
    """读取音频时长并为幻灯片设置自动翻页时间。

    停留时间 = 音频时长 + ``config.auto_advance_delay``。

    当音频文件缺失或时长解析失败时打印警告并跳过该页自动翻页，
    不影响整体处理流程（优雅降级）。警告同时经 ``on_progress`` 上报
    （作为 SETTING_TRANSITION 阶段的 message），服务端客户端可见。

    Args:
        slide: 已嵌入音频的目标幻灯片。
        audio_file: 对应的音频文件路径。
        config: 配音处理配置（读取 auto_advance_delay）。
        idx: 幻灯片页码（用于日志输出）。
        total: 幻灯片总页数（用于百分比计算）。
        start: 处理起始单调时间戳（用于 ETA 计算）。
        on_progress: 进度回调；为 None 时走打印输出。
    """
    try:
        duration = get_audio_duration(audio_file)
    except (OSError, ValueError) as exc:
        msg = f"   ⚠️ 第{idx}页无法读取音频时长，跳过自动翻页：{exc}"
        _log(on_progress, msg)
        _emit_progress(
            on_progress, start, STAGE_SETTING_TRANSITION, idx, total, msg
        )
        return

    delay = duration + config.auto_advance_delay
    try:
        set_advance_after_time(slide, delay)
    except (OSError, ValueError) as exc:
        msg = f"   ⚠️ 第{idx}页无法设置自动翻页：{exc}"
        _log(on_progress, msg)
        _emit_progress(
            on_progress, start, STAGE_SETTING_TRANSITION, idx, total, msg
        )
        return

    msg = (
        f"   ⏱️ 第{idx}页自动翻页：音频 {duration:.1f}s "
        f"+ 缓冲 {config.auto_advance_delay}s = {delay:.1f}s"
    )
    _log(on_progress, msg)
    _emit_progress(
        on_progress, start, STAGE_SETTING_TRANSITION, idx, total, msg
    )


async def process_slides(
    prs: Presentation,
    config: PTSpeechConfig,
    on_progress: "ProgressCallback | None" = None,
) -> None:
    """处理演示文稿中的所有幻灯片，生成并嵌入配音。

    按顺序遍历每张幻灯片：
    - 无备注：打印提示并跳过。
    - 有备注：调用 TTS 合成音频，成功后嵌入到幻灯片；若启用
      ``auto_advance``，则按「音频时长 + n 秒」设置该页自动翻页。
    处理完成后保存输出文件。

    中间音频文件存放在临时目录中，由 :func:`_temp_audio_workspace` 统一管理
    生命周期：进入 ``with`` 块时创建，退出时（无论成功或异常）自动清理，
    因此本函数不再需要手动 ``try/finally`` 删除临时文件。

    Args:
        prs: python-pptx 已打开的 Presentation 对象。
        config: 配音处理配置对象。
        on_progress: 可选进度回调。为 None 时（默认）保持与历史版本一致的
            ``print`` 输出；提供回调时改用结构化事件上报进度，原打印静默。

    Raises:
        EdgeTTSException: 当 TTS 合成流程中发生异常时向上抛出。
        OSError: 当输出 PPT 文件无法保存，或临时目录创建失败时向上抛出。
    """
    start = time.monotonic()
    total = len(prs.slides)

    _emit_progress(on_progress, start, STAGE_VALIDATING, 0, total, "校验配置")

    # 临时目录由上下文管理器统一管理：进入时创建，退出时自动清理。
    with _temp_audio_workspace(config) as temp_dir:
        for idx, slide in enumerate(prs.slides, start=1):
            _emit_progress(
                on_progress, start, STAGE_READING_NOTES, idx, total, "读取备注"
            )
            note_text = read_notes_text(slide)

            if not note_text:
                msg = f"【第{idx}页】无备注，跳过配音"
                _log(on_progress, msg)
                _emit_progress(
                    on_progress, start, STAGE_READING_NOTES, idx, total, msg
                )
                continue

            # 打印预览文本（超长截断）以便用户了解处理进度
            preview = note_text[:30] + "..." if len(note_text) > 30 else note_text
            msg = f"【第{idx}页】生成语音：{preview}"
            _log(on_progress, msg)
            _emit_progress(
                on_progress, start, STAGE_SYNTHESIZING, idx, total, msg
            )

            audio_file = temp_dir / f"slide_{idx}.mp3"
            success = await text_to_mp3(
                note_text,
                audio_file,
                voice_name=config.voice_name,
                speech_rate=config.speech_rate,
            )

            if success:
                _emit_progress(
                    on_progress, start, STAGE_EMBEDDING, idx, total, "嵌入音频"
                )
                embed_audio_autoplay(
                    slide,
                    audio_file,
                    icon_offset=config.audio_icon_offset,
                    icon_size=config.audio_icon_size,
                )

                if config.auto_advance:
                    _apply_auto_advance(
                        slide,
                        audio_file,
                        config,
                        idx,
                        total=total,
                        start=start,
                        on_progress=on_progress,
                    )

        _emit_progress(on_progress, start, STAGE_SAVING, 0, total, "保存输出文件")
        config.output_dir.mkdir(parents=True, exist_ok=True)
        prs.save(str(config.output_path))
        msg = f"\n✅ 处理完成！输出文件：{config.output_path}"
        _log(on_progress, msg)
        _emit_progress(
            on_progress, start, STAGE_COMPLETED, 0, total, "处理完成"
        )


async def speak_ppt_notes(
    config: Optional[PTSpeechConfig] = None,
    on_progress: "ProgressCallback | None" = None,
) -> None:
    """PPT 配音处理的顶层入口函数。

    典型用法：

    >>> import asyncio
    >>> from pathlib import Path
    >>> from ppt_speech import PTSpeechConfig, speak_ppt_notes
    >>> config = PTSpeechConfig(
    ...     input_dir=Path("data"),
    ...     output_dir=Path("data"),
    ...     voice_name="zh-CN-XiaoxiaoNeural",
    ...     speech_rate="+0%",
    ... )
    >>> asyncio.run(speak_ppt_notes(config))  # doctest: +SKIP

    Args:
        config: 配音处理配置；若为 None 则使用默认配置。
        on_progress: 可选进度回调，透传给 :func:`process_slides`。

    Raises:
        ValueError: 当配置参数（语音名称、语速等）不合法时。
        FileNotFoundError: 当输入 PPT 文件不存在时。
        EdgeTTSException: 当 TTS 服务请求失败时。
        OSError: 当文件读写或目录操作失败时。
    """
    if config is None:
        config = PTSpeechConfig(voice_name="zh-HK-WanLungNeural")

    config.validate()

    prs = Presentation(str(config.input_path))
    if on_progress is None:
        # 保持对 process_slides 的精确 2 参调用，兼容既有测试断言。
        await process_slides(prs, config)
    else:
        await process_slides(prs, config, on_progress=on_progress)